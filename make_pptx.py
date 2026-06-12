"""
Generate MineSafe AI System — HITSZ 2026 PowerPoint presentation.
Run: python3 make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0b, 0x0e, 0x1a)
CARD     = RGBColor(0x19, 0x20, 0x30)
CARD2    = RGBColor(0x14, 0x19, 0x29)
BLUE     = RGBColor(0x3b, 0x82, 0xf6)
CYAN     = RGBColor(0x06, 0xb6, 0xd4)
GREEN    = RGBColor(0x10, 0xb9, 0x81)
RED      = RGBColor(0xef, 0x44, 0x44)
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)
WHITE    = RGBColor(0xf1, 0xf5, 0xf9)
GREY     = RGBColor(0x94, 0xa3, 0xb8)
DIM      = RGBColor(0x47, 0x55, 0x69)
BORDER   = RGBColor(0x1e, 0x29, 0x3b)

W  = Inches(13.33)   # slide width  (16:9 widescreen)
H  = Inches(7.5)     # slide height

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank

def add_slide(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, BG)
    return slide

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, x, y, w, h, text, size=Pt(11), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf_box

def add_para(tf, text, size=Pt(10), bold=False, color=GREY,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(2)):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def accent_bar(slide, color=BLUE, thickness=Inches(0.055)):
    rect(slide, 0, 0, W, thickness, fill=color)

def slide_number_label(slide, n, x=Inches(0.45), y=Inches(0.8)):
    txbox(slide, x, y, Inches(1.2), Inches(0.55),
          f"{n:02d}", size=Pt(28), bold=True, color=BORDER)

def section_label(slide, text, x, y, w=Inches(5), color=BLUE):
    txbox(slide, x, y, w, Inches(0.3),
          text.upper(), size=Pt(7), bold=True, color=color)

def divider_line(slide, x, y, w, color=BORDER, thickness=Pt(0.5)):
    from pptx.util import Pt
    ln = slide.shapes.add_shape(1, x, y, w, Pt(1))
    ln.fill.background()
    ln.line.color.rgb = color
    ln.line.width = thickness

def stat_card(slide, x, y, w, h, value, label, accent=BLUE):
    rect(slide, x, y, w, h, fill=CARD, line=BORDER)
    rect(slide, x, y, w, Inches(0.055), fill=accent)
    txbox(slide, x, y + Inches(0.18), w, Inches(0.7),
          value, size=Pt(32), bold=True, color=accent, align=PP_ALIGN.CENTER)
    txbox(slide, x, y + Inches(0.85), w, Inches(0.5),
          label, size=Pt(8), color=GREY, align=PP_ALIGN.CENTER, wrap=True)

def bullet_list(slide, items, x, y, w, h, icon='▸', size=Pt(9.5)):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        # strip JS-style tags like <1%>
        clean = item.replace('<', '').replace('>', '')
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{icon}  {clean}"
        run.font.size = size
        run.font.color.rgb = GREY
    return tf_box

# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)

    # radial-ish overlay at top-right
    rect(sl, Inches(6), 0, Inches(7.33), Inches(4.5), fill=RGBColor(0x0e, 0x20, 0x40))

    txbox(sl, Inches(0.9), Inches(0.4), Inches(7), Inches(0.35),
          "AI FOR SUSTAINABLE DEVELOPMENT", size=Pt(8), bold=True, color=BLUE)

    txbox(sl, Inches(0.9), Inches(0.85), Inches(11), Inches(1.6),
          s['title'], size=Pt(52), bold=True, color=WHITE)

    txbox(sl, Inches(0.9), Inches(2.5), Inches(9), Inches(0.5),
          s['subtitle'], size=Pt(14), color=GREY, italic=True)

    divider_line(sl, Inches(0.9), Inches(3.1), Inches(2), color=BLUE)

    txbox(sl, Inches(0.9), Inches(3.25), Inches(6), Inches(0.32),
          s.get('institution',''), size=Pt(9), bold=True, color=GREY)

    members = s.get('members', [])
    mx = Inches(0.9)
    for i, m in enumerate(members):
        rect(sl, mx + i * Inches(2.55), Inches(3.75), Inches(2.4), Inches(0.42),
             fill=CARD, line=BORDER)
        txbox(sl, mx + i * Inches(2.55) + Inches(0.1), Inches(3.8),
              Inches(2.2), Inches(0.32), m, size=Pt(9), color=GREY)

    txbox(sl, Inches(0.9), Inches(4.32), Inches(4), Inches(0.3),
          s.get('year',''), size=Pt(9), color=DIM)

    if s.get('url'):
        txbox(sl, Inches(0.9), Inches(4.72), Inches(5), Inches(0.28),
              "Live Dashboard:", size=Pt(8), bold=True, color=DIM)
        txbox(sl, Inches(0.9), Inches(5.0), Inches(6), Inches(0.32),
              s['url'], size=Pt(10), bold=True, color=BLUE)

    txbox(sl, Inches(0.7), Inches(6.8), Inches(12), Inches(0.35),
          f"Slide {num}", size=Pt(7), color=DIM)
    return sl


def slide_team(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    txbox(sl, Inches(1.8), Inches(0.75), Inches(10), Inches(0.65),
          s['title'], size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(1.8), Inches(1.35), Inches(10), Inches(0.32),
          s.get('subtitle',''), size=Pt(10), color=GREY,
          align=PP_ALIGN.CENTER, italic=True)

    members = s.get('members', [])
    card_w = Inches(2.85)
    card_h = Inches(3.6)
    gap    = Inches(0.28)
    total  = len(members) * card_w + (len(members)-1) * gap
    start_x = (W - total) / 2

    roles_colors = [BLUE, CYAN, GREEN, AMBER]

    for i, m in enumerate(members):
        cx = start_x + i * (card_w + gap)
        cy = Inches(1.85)
        rect(sl, cx, cy, card_w, card_h, fill=CARD, line=BORDER)
        rect(sl, cx, cy, card_w, Inches(0.045), fill=roles_colors[i % 4])

        # avatar circle placeholder
        av_x = cx + (card_w - Inches(1.1)) / 2
        av_y = cy + Inches(0.3)
        av = sl.shapes.add_shape(9, av_x, av_y, Inches(1.1), Inches(1.1))  # oval
        av.fill.solid()
        av.fill.fore_color.rgb = roles_colors[i % 4]
        av.line.fill.background()
        initials = ''.join(w[0] for w in m['name'].split()[:2])
        txbox(sl, av_x, av_y + Inches(0.3), Inches(1.1), Inches(0.5),
              initials, size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txbox(sl, cx + Inches(0.1), cy + Inches(1.6), card_w - Inches(0.2), Inches(0.45),
              m['name'], size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, cx + Inches(0.1), cy + Inches(2.05), card_w - Inches(0.2), Inches(0.9),
              m['role'], size=Pt(8), color=GREY, align=PP_ALIGN.CENTER, wrap=True)

    txbox(sl, 0, Inches(5.65), W, Inches(0.3),
          "HITSZ — Harbin Institute of Technology, Shenzhen",
          size=Pt(8), color=DIM, align=PP_ALIGN.CENTER)
    return sl


def slide_content(prs, s, num, accent=BLUE):
    sl = add_slide(prs)
    accent_bar(sl, accent)
    slide_number_label(sl, num)

    lw = Inches(5.8)
    lx = Inches(0.55)
    ly = Inches(0.75)

    txbox(sl, lx, ly, lw, Inches(0.55),
          s['title'], size=Pt(22), bold=True, color=WHITE)
    txbox(sl, lx, ly + Inches(0.6), lw, Inches(0.35),
          s.get('subtitle',''), size=Pt(9.5), color=GREY, italic=True)

    divider_line(sl, lx, ly + Inches(1.05), Inches(1.5), color=accent)

    items = s.get('content', [])
    bullet_list(sl, items, lx, ly + Inches(1.2), lw, Inches(3.5))

    # right panel
    rx = Inches(6.8)
    ry = Inches(0.75)
    rw = Inches(6.1)

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, rx, ry, rw, Inches(1.55),
                  stat.get('value',''), stat.get('label',''), accent=accent)

    # sources
    sources = s.get('sources', [])
    if sources:
        sy = Inches(6.5)
        txbox(sl, lx, sy, Inches(12), Inches(0.25),
              "Sources: " + "  ·  ".join(f"{s['label']} ({s['year']})" for s in sources),
              size=Pt(6.5), color=DIM, italic=True)
    return sl


def slide_problem(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    lx, ly, lw = Inches(0.55), Inches(0.75), Inches(5.8)

    txbox(sl, lx, ly, lw, Inches(0.55),
          s['title'], size=Pt(22), bold=True, color=WHITE)
    txbox(sl, lx, ly + Inches(0.6), lw, Inches(0.32),
          s.get('subtitle',''), size=Pt(9.5), color=GREY, italic=True)

    # highlighted question
    qy = ly + Inches(1.1)
    rect(sl, lx, qy, lw, Inches(0.9), fill=RGBColor(0x0f, 0x1e, 0x3a), line=BLUE, line_w=Pt(0.5))
    rect(sl, lx, qy, Inches(0.06), Inches(0.9), fill=BLUE)
    txbox(sl, lx + Inches(0.15), qy + Inches(0.08), lw - Inches(0.25), Inches(0.72),
          s.get('question',''), size=Pt(8.5), color=WHITE, italic=True, wrap=True)

    items = s.get('points', [])
    bullet_list(sl, items, lx, qy + Inches(1.05), lw, Inches(3.2))

    rx, ry, rw = Inches(6.8), Inches(0.75), Inches(6.1)
    stat = s.get('stat', {})
    if stat:
        stat_card(sl, rx, ry, rw, Inches(1.55),
                  stat.get('value',''), stat.get('label',''))

    # flow steps
    steps = ['Detect risk early', 'Explain the cause', 'Recommend action', 'Prevent the incident']
    for i, step in enumerate(steps):
        sy = ry + Inches(1.8) + i * Inches(0.85)
        rect(sl, rx, sy, rw, Inches(0.7), fill=CARD, line=BORDER)
        rect(sl, rx, sy + Inches(0.28), Inches(0.07), Inches(0.14), fill=BLUE)
        txbox(sl, rx + Inches(0.22), sy + Inches(0.2), rw - Inches(0.35), Inches(0.35),
              step, size=Pt(9.5), color=GREY)
    return sl


def slide_comparison(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    txbox(sl, Inches(1.8), Inches(0.75), Inches(9.5), Inches(0.55),
          s['title'], size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(1.8), Inches(1.3), Inches(9.5), Inches(0.32),
          s.get('subtitle',''), size=Pt(9.5), color=GREY,
          align=PP_ALIGN.CENTER, italic=True)

    col_w = Inches(5.6)
    col_h = Inches(4.2)
    lx, rx = Inches(0.5), Inches(7.0)
    cy = Inches(1.75)

    # old column
    rect(sl, lx, cy, col_w, col_h, fill=CARD, line=BORDER)
    rect(sl, lx, cy, col_w, Inches(0.045), fill=RED)
    txbox(sl, lx + Inches(0.2), cy + Inches(0.1), col_w - Inches(0.4), Inches(0.35),
          "OLD WAY", size=Pt(8), bold=True, color=RED)
    for i, item in enumerate(s.get('old', [])):
        iy = cy + Inches(0.55) + i * Inches(0.58)
        txbox(sl, lx + Inches(0.2), iy, col_w - Inches(0.35), Inches(0.5),
              f"✗  {item}", size=Pt(9), color=GREY, wrap=True)

    # arrow
    txbox(sl, Inches(6.35), cy + Inches(1.9), Inches(0.6), Inches(0.5),
          "→", size=Pt(22), color=DIM, align=PP_ALIGN.CENTER)

    # new column
    rect(sl, rx, cy, col_w, col_h, fill=CARD, line=BORDER)
    rect(sl, rx, cy, col_w, Inches(0.045), fill=BLUE)
    txbox(sl, rx + Inches(0.2), cy + Inches(0.1), col_w - Inches(0.4), Inches(0.35),
          "MINESAFE WAY", size=Pt(8), bold=True, color=BLUE)
    for i, item in enumerate(s.get('new', [])):
        iy = cy + Inches(0.55) + i * Inches(0.58)
        txbox(sl, rx + Inches(0.2), iy, col_w - Inches(0.35), Inches(0.5),
              f"✓  {item}", size=Pt(9), color=GREY, wrap=True)

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, Inches(4.4), Inches(6.2), Inches(4.5), Inches(1.1),
                  stat.get('value',''), stat.get('label',''))
    return sl


def slide_solution(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    txbox(sl, Inches(1.5), Inches(0.75), Inches(10), Inches(0.55),
          s['title'], size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(1.5), Inches(1.3), Inches(10), Inches(0.32),
          s.get('subtitle',''), size=Pt(9.5), color=GREY,
          align=PP_ALIGN.CENTER, italic=True)

    col_w = Inches(4.9)
    col_h = Inches(4.5)
    cy = Inches(1.8)
    lx, rx = Inches(0.5), Inches(7.9)

    rect(sl, lx, cy, col_w, col_h, fill=CARD, line=BORDER)
    txbox(sl, lx + Inches(0.15), cy + Inches(0.1), col_w - Inches(0.3), Inches(0.28),
          "CONNECTS", size=Pt(7), bold=True, color=DIM)
    for i, item in enumerate(s.get('inputs', [])):
        iy = cy + Inches(0.45) + i * Inches(0.49)
        rect(sl, lx + Inches(0.15), iy, col_w - Inches(0.3), Inches(0.38),
             fill=RGBColor(0x1a, 0x2a, 0x45), line=RGBColor(0x2a, 0x45, 0x70))
        # strip emoji for cleaner display
        clean = item.strip()
        txbox(sl, lx + Inches(0.25), iy + Inches(0.05), col_w - Inches(0.5), Inches(0.3),
              clean, size=Pt(8.5), color=GREY)

    txbox(sl, Inches(5.55), cy + Inches(2.1), Inches(0.8), Inches(0.5),
          "→", size=Pt(24), color=DIM, align=PP_ALIGN.CENTER)

    rect(sl, rx, cy, col_w, col_h, fill=CARD, line=BORDER)
    txbox(sl, rx + Inches(0.15), cy + Inches(0.1), col_w - Inches(0.3), Inches(0.28),
          "OUTPUTS", size=Pt(7), bold=True, color=DIM)
    for i, item in enumerate(s.get('outputs', [])):
        iy = cy + Inches(0.45) + i * Inches(0.65)
        rect(sl, rx + Inches(0.15), iy, col_w - Inches(0.3), Inches(0.5),
             fill=RGBColor(0x0d, 0x22, 0x1a), line=RGBColor(0x10, 0x40, 0x2a))
        txbox(sl, rx + Inches(0.25), iy + Inches(0.07), col_w - Inches(0.5), Inches(0.38),
              item, size=Pt(8.5), color=GREY)

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, Inches(4.4), Inches(6.45), Inches(4.5), Inches(0.85),
                  stat.get('value',''), stat.get('label',''))
    return sl


def slide_agents(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    lx, ly, lw = Inches(0.55), Inches(0.75), Inches(5.5)
    txbox(sl, lx, ly, lw, Inches(0.55),
          s['title'], size=Pt(20), bold=True, color=WHITE)
    txbox(sl, lx, ly + Inches(0.58), lw, Inches(0.32),
          s.get('subtitle',''), size=Pt(9), color=GREY, italic=True)
    bullet_list(sl, s.get('content',[]), lx, ly + Inches(1.0), lw, Inches(2.0), size=Pt(8.5))

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, lx, ly + Inches(3.1), lw, Inches(1.1),
                  stat.get('value',''), stat.get('label',''))

    # agents right panel
    rx = Inches(6.5)
    ry = Inches(0.6)
    rw = Inches(6.4)

    txbox(sl, rx, ry, rw, Inches(0.28),
          "LIVE AGENT ACTIVATIONS", size=Pt(7), bold=True, color=DIM)

    status_colors = {'critical': RED, 'warning': AMBER, 'normal': GREEN}
    for i, ag in enumerate(s.get('agents', [])):
        ay = ry + Inches(0.38) + i * Inches(1.2)
        sc = status_colors.get(ag.get('status','normal'), GREEN)
        rect(sl, rx, ay, rw, Inches(1.05), fill=CARD, line=BORDER)
        rect(sl, rx, ay, Inches(0.055), Inches(1.05), fill=sc)
        txbox(sl, rx + Inches(0.3), ay + Inches(0.08), rw - Inches(0.8), Inches(0.32),
              ag.get('name',''), size=Pt(9.5), bold=True, color=sc)
        case = ag.get('case','')
        if len(case) > 110:
            case = case[:107] + '…'
        txbox(sl, rx + Inches(0.3), ay + Inches(0.42), rw - Inches(0.5), Inches(0.55),
              case, size=Pt(7.5), color=GREY, wrap=True)
        dot = sl.shapes.add_shape(9, rx + rw - Inches(0.3), ay + Inches(0.42),
                                  Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = sc; dot.line.fill.background()
    return sl


def slide_training(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, GREEN)
    slide_number_label(sl, num)

    lx, ly, lw = Inches(0.55), Inches(0.75), Inches(5.5)
    txbox(sl, lx, ly, lw, Inches(0.55),
          s['title'], size=Pt(20), bold=True, color=WHITE)
    txbox(sl, lx, ly + Inches(0.58), lw, Inches(0.32),
          s.get('subtitle',''), size=Pt(9), color=GREY, italic=True)
    txbox(sl, lx, ly + Inches(1.0), lw, Inches(0.85),
          s.get('intro',''), size=Pt(8.5), color=GREY, italic=True, wrap=True)

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, lx, ly + Inches(2.0), lw, Inches(1.1),
                  stat.get('value',''), stat.get('label',''), accent=GREEN)

    rx = Inches(6.5)
    ry = Inches(0.6)
    rw = Inches(6.4)
    txbox(sl, rx, ry, rw, Inches(0.28),
          "AIESD PRINCIPLES APPLIED", size=Pt(7), bold=True, color=DIM)

    principles = s.get('principles', [])
    cols = 2
    card_w = (rw - Inches(0.15)) / cols
    card_h = Inches(1.15)
    gap = Inches(0.12)

    for i, p in enumerate(principles):
        col = i % cols
        row = i // cols
        px = rx + col * (card_w + gap)
        py = ry + Inches(0.38) + row * (card_h + gap)
        rect(sl, px, py, card_w, card_h, fill=CARD, line=BORDER)
        rect(sl, px, py, Inches(0.045), card_h, fill=GREEN)
        txbox(sl, px + Inches(0.14), py + Inches(0.1), card_w - Inches(0.18), Inches(0.32),
              f"{p['icon']}  {p['title']}", size=Pt(9), bold=True, color=WHITE)
        txbox(sl, px + Inches(0.14), py + Inches(0.46), card_w - Inches(0.18), Inches(0.62),
              p['desc'], size=Pt(7.5), color=GREY, wrap=True)
    return sl


def slide_visits(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    txbox(sl, Inches(1.5), Inches(0.75), Inches(10), Inches(0.52),
          s['title'], size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(1.5), Inches(1.25), Inches(10), Inches(0.32),
          s.get('subtitle',''), size=Pt(9.5), color=GREY,
          align=PP_ALIGN.CENTER, italic=True)

    visits = s.get('visits', [])
    n = len(visits)
    margin = Inches(0.35)
    gap    = Inches(0.18)
    total_gap = gap * (n - 1)
    card_w = (W - 2 * margin - total_gap) / n
    card_h = Inches(4.6)
    cy     = Inches(1.72)

    vc_colors = [CYAN, BLUE, GREEN, RGBColor(0xa8, 0x55, 0xf7), AMBER]

    for i, v in enumerate(visits):
        vx = margin + i * (card_w + gap)
        vc = vc_colors[i % len(vc_colors)]
        rect(sl, vx, cy, card_w, card_h, fill=CARD, line=BORDER)
        rect(sl, vx, cy, card_w, Inches(0.05), fill=vc)
        txbox(sl, vx + Inches(0.12), cy + Inches(0.1), card_w - Inches(0.2), Inches(0.32),
              v.get('company',''), size=Pt(9.5), bold=True, color=vc)
        txbox(sl, vx + Inches(0.12), cy + Inches(0.42), card_w - Inches(0.2), Inches(0.28),
              v.get('tagline',''), size=Pt(7), color=DIM)
        divider_line(sl, vx + Inches(0.12), cy + Inches(0.8), card_w - Inches(0.24), color=BORDER)
        txbox(sl, vx + Inches(0.1), cy + Inches(0.92), card_w - Inches(0.2), Inches(0.22),
              "LESSON", size=Pt(6.5), bold=True, color=DIM)
        txbox(sl, vx + Inches(0.1), cy + Inches(1.12), card_w - Inches(0.2), Inches(1.0),
              v.get('lesson',''), size=Pt(7.5), color=GREY, wrap=True)
        divider_line(sl, vx + Inches(0.12), cy + Inches(2.2), card_w - Inches(0.24), color=BORDER)
        txbox(sl, vx + Inches(0.1), cy + Inches(2.32), card_w - Inches(0.2), Inches(0.22),
              "MINESAFE →", size=Pt(6.5), bold=True, color=CYAN)
        txbox(sl, vx + Inches(0.1), cy + Inches(2.55), card_w - Inches(0.2), Inches(1.0),
              v.get('connection',''), size=Pt(7.5), color=CYAN, wrap=True)
    return sl


def slide_dashboard(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)
    slide_number_label(sl, num)

    lx, ly, lw = Inches(0.55), Inches(0.75), Inches(5.5)
    txbox(sl, lx, ly, lw, Inches(0.55),
          s['title'], size=Pt(20), bold=True, color=WHITE)
    txbox(sl, lx, ly + Inches(0.58), lw, Inches(0.32),
          s.get('subtitle',''), size=Pt(9), color=GREY, italic=True)
    bullet_list(sl, s.get('content',[]), lx, ly + Inches(1.0), lw, Inches(2.5))

    stat = s.get('stat', {})
    if stat:
        stat_card(sl, lx, ly + Inches(3.6), lw, Inches(1.2),
                  stat.get('value',''), stat.get('label',''))

    # right: live dashboard placeholder
    rx, ry, rw, rh = Inches(6.5), Inches(0.6), Inches(6.4), Inches(5.8)
    rect(sl, rx, ry, rw, rh, fill=CARD, line=BORDER)
    rect(sl, rx, ry, rw, Inches(0.045), fill=GREEN)

    # live dot label
    dot = sl.shapes.add_shape(9, rx + Inches(0.2), ry + Inches(0.12),
                               Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()
    txbox(sl, rx + Inches(0.42), ry + Inches(0.08), Inches(2.5), Inches(0.25),
          "LIVE DASHBOARD PREVIEW", size=Pt(6.5), bold=True, color=GREEN)

    route = s.get('route','')
    if route:
        txbox(sl, rx + Inches(0.2), ry + Inches(2.5), rw - Inches(0.4), Inches(0.4),
              f"minesafe.iraw.one/#{route}", size=Pt(9), color=BLUE, align=PP_ALIGN.CENTER)

    txbox(sl, rx + Inches(0.2), ry + Inches(3.0), rw - Inches(0.4), Inches(0.3),
          "Visit the live dashboard for a real-time demo",
          size=Pt(8), color=GREY, align=PP_ALIGN.CENTER, italic=True)
    return sl


def slide_takeaway(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, GREEN)
    fill_bg(sl, RGBColor(0x09, 0x0e, 0x16))

    txbox(sl, 0, Inches(0.7), W, Inches(0.65),
          s.get('icon','✨') + "  Takeaway",
          size=Pt(14), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    txbox(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5),
          s.get('title', 'Takeaway'),
          size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    divider_line(sl, Inches(5.2), Inches(3.2), Inches(2.9), color=GREEN)

    txbox(sl, Inches(1.2), Inches(3.4), Inches(10.9), Inches(1.5),
          s.get('statement',''),
          size=Pt(11), color=GREY, align=PP_ALIGN.CENTER, wrap=True)

    txbox(sl, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.9),
          s.get('closing',''),
          size=Pt(11), bold=True, color=GREEN,
          align=PP_ALIGN.CENTER, italic=True, wrap=True)
    return sl


def slide_thankyou(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, GREEN)
    fill_bg(sl, RGBColor(0x09, 0x0e, 0x16))

    txbox(sl, 0, Inches(1.0), W, Inches(0.55),
          s.get('icon','🎓') + "  " + s.get('institution',''),
          size=Pt(11), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    txbox(sl, 0, Inches(1.7), W, Inches(1.6),
          s.get('title','Thank You'),
          size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    divider_line(sl, Inches(4.5), Inches(3.55), Inches(4.3), color=GREEN)

    txbox(sl, 0, Inches(3.75), W, Inches(0.5),
          s.get('subtitle',''),
          size=Pt(13), color=GREY, align=PP_ALIGN.CENTER)

    members = s.get('members', [])
    card_w = Inches(2.7)
    total  = len(members) * card_w + (len(members) - 1) * Inches(0.2)
    mx = (W - total) / 2
    for i, m in enumerate(members):
        rect(sl, mx + i * (card_w + Inches(0.2)), Inches(4.45),
             card_w, Inches(0.45), fill=CARD, line=BORDER)
        txbox(sl, mx + i * (card_w + Inches(0.2)) + Inches(0.1), Inches(4.5),
              card_w - Inches(0.2), Inches(0.35), m, size=Pt(9), color=GREY,
              align=PP_ALIGN.CENTER)

    txbox(sl, 0, Inches(5.2), W, Inches(0.32),
          "MineSafe AI System · AI-Driven Mining Safety & Sustainability",
          size=Pt(8), color=DIM, align=PP_ALIGN.CENTER)
    return sl


def slide_video_thankyou(prs, s, num):
    sl = add_slide(prs)
    accent_bar(sl, BLUE)

    txbox(sl, 0, Inches(0.5), W, Inches(0.28),
          "CLOSING CEREMONY", size=Pt(7), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    txbox(sl, 0, Inches(0.85), W, Inches(1.2),
          s.get('title','Thank You'),
          size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txbox(sl, 0, Inches(2.0), W, Inches(0.4),
          s.get('subtitle',''),
          size=Pt(11), color=GREY, align=PP_ALIGN.CENTER, italic=True)

    txbox(sl, 0, Inches(2.45), W, Inches(0.32),
          s.get('closing',''),
          size=Pt(8.5), color=DIM, align=PP_ALIGN.CENTER)

    # video placeholder card
    vx, vy, vw, vh = Inches(2.3), Inches(2.95), Inches(8.7), Inches(3.55)
    rect(sl, vx, vy, vw, vh, fill=RGBColor(0x04, 0x09, 0x12), line=BLUE, line_w=Pt(1.2))

    # play button circle
    pw = Inches(1.0)
    px = vx + (vw - pw) / 2
    py = vy + (vh - pw) / 2 - Inches(0.3)
    circle = sl.shapes.add_shape(9, px, py, pw, pw)
    circle.fill.solid(); circle.fill.fore_color.rgb = BLUE; circle.line.fill.background()
    txbox(sl, px + Inches(0.08), py + Inches(0.22), pw - Inches(0.1), Inches(0.55),
          "▶", size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)

    txbox(sl, vx, vy + vh - Inches(0.85), vw, Inches(0.32),
          "Watch Team Memory Video", size=Pt(9), bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, vx, vy + vh - Inches(0.52), vw, Inches(0.28),
          "♪ See You Again", size=Pt(8), color=GREY,
          align=PP_ALIGN.CENTER, italic=True)

    txbox(sl, 0, Inches(6.75), W, Inches(0.28),
          s.get('footer', 'MineSafe AI System · HITSZ 2026'),
          size=Pt(7), color=DIM, align=PP_ALIGN.CENTER)

    # download note
    dl = s.get('downloadSrc','')
    if dl:
        txbox(sl, 0, Inches(6.5), W, Inches(0.28),
              "↓  Full video download available at the live dashboard",
              size=Pt(7), color=BLUE, align=PP_ALIGN.CENTER)
    return sl


# ── Slide data ────────────────────────────────────────────────────────────────

slides_data = [
  { 'id':0, 'type':'title',
    'title':'MineSafe AI System',
    'subtitle':'AI-Driven Mining Safety & Sustainability Command Center',
    'institution':'Harbin Institute of Technology, Shenzhen (HITSZ)',
    'course':'AI for Sustainable Development',
    'year':'2026',
    'url':'https://minesafe.iraw.one',
    'members':['Harta Deddy Irawan','Sashi Kumar','Isabella Gbefa','Shania Clauren'],
  },
  { 'id':1, 'type':'team',
    'title':'Meet the Team',
    'subtitle':'AI for Sustainable Development — HITSZ 2026',
    'members':[
      { 'id':'harta',    'name':'Harta Deddy Irawan', 'role':'Lead Developer & System Architecture' },
      { 'id':'sashi',    'name':'Sashi Kumar',         'role':'Worker Health & Medical Safety' },
      { 'id':'isabella', 'name':'Isabella Gbefa',      'role':'Health and Safety, ESG Officer' },
      { 'id':'shania',   'name':'Shania Clauren',      'role':'Finance & Business Value' },
    ],
  },
  { 'id':2, 'type':'content',
    'title':'The Mining Safety Challenge',
    'subtitle':'Why AI is essential — and why mining is the right environment to demonstrate it',
    'content':[
      'Mining employs <1% of workers yet contributes a disproportionate share of fatal work accidents (ILO, 2023)',
      'Workers face simultaneous exposure: gas, heat, dust, noise, slope instability, and fatigue in one shift',
      'Safety data is fragmented across sensors, wearables, vehicles, equipment, and drones — all disconnected',
      'Environmental incidents cost $6B+ annually in fines, remediation, and reputational damage',
      'ESG disclosure mandatory in 60+ jurisdictions; investors demand verified sustainability data',
      'Mining is an ideal environment to demonstrate AI for safety and sustainable development',
    ],
    'stat':{ 'value':'5,000+', 'label':'Estimated mining fatalities / year (ILO, 2023)' },
    'sources':[
      { 'label':'ILO — Safety and Health at Work Statistics', 'year':'2023' },
      { 'label':'ICMM — Safety Performance Benchmarking', 'year':'2024' },
      { 'label':'Deloitte — Tracking the Trends 2024: Mining & Metals', 'year':'2024' },
    ],
  },
  { 'id':3, 'type':'problem',
    'title':'The Problem We Want to Solve',
    'subtitle':'One question drives the entire system',
    'question':'"How can mining operations detect and prevent safety, health, and environmental risks before they become incidents?"',
    'points':[
      'Mining risks are dynamic and distributed across hundreds of separate data sources',
      'Safety data is fragmented: sensors, workers, vehicles, equipment, drones, and environment are all disconnected',
      'Traditional safety monitoring is reactive — operators learn about incidents after they happen',
      'Operators need early warning, explainable risk predictions, and recommended actions',
      'The goal is prevention, not only incident response',
    ],
    'stat':{ 'value':'72%', 'label':'Of mining incidents are preventable with earlier intervention' },
  },
  { 'id':5, 'type':'comparison',
    'title':'From Reactive Monitoring to Predictive Prevention',
    'subtitle':'How MineSafe changes the operational model',
    'old':[
      'Separate, disconnected dashboards',
      'Manual inspection and site walkthrough',
      'Periodic safety reports',
      'Delayed response after an incident',
      'Compliance reviewed after the fact',
      'Reactive incident management',
    ],
    'new':[
      'Continuous multi-source sensing across all zones',
      'AI risk prediction with 24-hour horizon',
      'Explainable alerts with confidence scores',
      'Real-time command center for all operators',
      'Human-in-the-loop decision support',
      'Prevention-focused recommended actions',
    ],
    'stat':{ 'value':'89%', 'label':'Maximum prediction confidence — gas risk model' },
  },
  { 'id':6, 'type':'solution',
    'title':'What Problem MineSafe Solves',
    'subtitle':'Connecting fragmented mining safety intelligence into one command center',
    'inputs':[
      'Gas & environmental sensors',
      'Worker health wearables',
      'Drones & visual inspection',
      'Fleet GPS & driver risk',
      'Equipment telemetry',
      'Noise exposure dosimeters',
      'Sustainability indicators',
      'Incident response workflow',
    ],
    'outputs':[
      'Early warning alerts',
      '24-hour risk forecast',
      'AI confidence scoring',
      'Explainable AI decisions',
      'Ranked recommended actions',
      'ESG & operational visibility',
    ],
    'stat':{ 'value':'8', 'label':'Safety domains connected in one command center' },
  },
  { 'id':8, 'type':'agents',
    'title':'AI Agents in the Analytics Layer',
    'subtitle':'Five autonomous agents — one unified safety intelligence',
    'content':[
      'Goal-driven programs that perceive, reason, and act without continuous human input',
      'Each agent monitors its domain, generates predictions with confidence scores, and triggers action workflows',
      'Agents collaborate: a gas spike from Agent 1 automatically informs Agent 5 to pre-calculate evacuation routes',
      'Every decision is explainable — reading, status, confidence %, and recommended action are shown',
    ],
    'stat':{ 'value':'5', 'label':'Autonomous AI Agents operating across the mine' },
    'agents':[
      { 'icon':'🔬', 'name':'Gas Safety Agent',
        'case':'Watches GS-11 CH4 8.4 ppm at Zone 8. Predicts threshold breach within 2 hours. Recommends evacuation + ventilation.',
        'status':'critical' },
      { 'icon':'💗', 'name':'Worker Welfare Agent',
        'case':'Flags W004 heat stress (fatigue 78) and W010 fall detected (no movement 90s). Dispatches rescue team.',
        'status':'critical' },
      { 'icon':'⚙️', 'name':'Equipment Maintenance Agent',
        'case':'EQ08 engine temp 103C — 10h before maintenance due. Predicts failure, schedules emergency stop.',
        'status':'warning' },
      { 'icon':'🚁', 'name':'Drone Commander Agent',
        'case':'Eagle-2 detects slope crack on Zone B north face (conf. 84%). Redirects Eagle-1 for confirmation pass.',
        'status':'warning' },
      { 'icon':'🌿', 'name':'Environmental Compliance Agent',
        'case':'Monitors GS-08 PM2.5 58.4 ug/m3 and tailings seepage risk. Auto-notifies regulator on threshold breach.',
        'status':'normal' },
    ],
  },
  { 'id':9, 'type':'training',
    'title':'Connection to AIESD Training',
    'subtitle':'Applying AI for Sustainable Development principles to a real industrial use case',
    'intro':'MineSafe applies AIESD learning by turning AI ethics and sustainable development principles into a practical safety technology concept.',
    'stat':{ 'value':'AIESD', 'label':'AI for Sustainable Development — HITSZ 2026' },
    'principles':[
      { 'icon':'🧑‍💼', 'title':'Human-centered AI',    'desc':'AI supports supervisors and safety teams instead of replacing them. Every decision stays with a qualified operator.' },
      { 'icon':'🔍',   'title':'Explainable AI (XAI)', 'desc':'Each prediction shows why it was triggered — factors, confidence score, data sources, and data quality.' },
      { 'icon':'⚖️',   'title':'Responsible AI',       'desc':'Confidence scores and data quality are always visible. No recommendation is issued without a clear reasoning chain.' },
      { 'icon':'🌍',   'title':'Sustainable Development','desc':'Worker safety, environmental monitoring, ESG indicators, and prevention of long-term occupational harm.' },
      { 'icon':'🏭',   'title':'Practical Implementation','desc':'Designed to demonstrate how AI can be deployed in real industrial environments with existing IoT infrastructure.' },
    ],
  },
  { 'id':9, 'type':'visits',
    'title':'Lessons from Industry Visits',
    'subtitle':'How each company visit shaped MineSafe\'s design',
    'visits':[
      { 'company':'Chipsea',          'tagline':'Smart IoT Chips',                'lesson':'Edge sensing and low-power chips enable continuous real-world monitoring at scale.', 'connection':'Gas sensors, noise sensors, environmental nodes, and equipment telemetry.' },
      { 'company':'Mindray',          'tagline':'Medical & Health Devices',       'lesson':'Reliable vital-sign monitoring supports timely safety decisions in demanding environments.', 'connection':'Worker heart rate, SpO2, body temp, fatigue, and heat stress monitoring.' },
      { 'company':'BYD',              'tagline':'Intelligent Vehicles & Energy',  'lesson':'Intelligent vehicles, batteries, and industrial-scale energy systems transform operations.', 'connection':'Fleet telemetry, vehicle status, driver risk, battery/fuel monitoring.' },
      { 'company':'Dobot',            'tagline':'Robotics & Automation',          'lesson':'Robotics automates inspection and reduces human exposure to hazardous environments.', 'connection':'Drone inspection, robotic inspection concept, and remote hazard detection.' },
      { 'company':'Fengnong Holdings','tagline':'Digital Agriculture & Sustainability','lesson':'Digital technology supports sustainable operations, environmental control, and resource efficiency.', 'connection':'ESG dashboard, dust, water, runoff, land rehabilitation, sustainability intelligence.' },
    ],
  },
  { 'id':11, 'type':'dashboard', 'title':'Command Center Overview',
    'subtitle':'Every alert. Every worker. Every risk. One screen.',
    'content':[
      'AI-generated aerial mine map with live overlays — drones, workers, risk zones, sensor readings',
      'Triaged alert panel: AI assigns severity, zone, and recommended action in < 1 second',
      '8 real-time KPIs: workers, drones, sensors, critical alerts, risk score, fatigue, PPE, evacuation',
      'Scenario simulation: trigger Gas Leak, Heat Stress, Slope Crack, Vehicle Proximity, or Sustainability Risk',
    ],
    'stat':{ 'value':'< 30s', 'label':'Alert-to-response latency' }, 'route':'/overview' },
  { 'id':12, 'type':'dashboard', 'title':'Drone Monitoring',
    'subtitle':'Aerial surveillance and slope crack detection in real time',
    'content':[
      '6 autonomous drones (Eagle-1 to Eagle-6) with live telemetry: battery, altitude, speed, zone',
      'Eagle-2 currently monitoring north slope crack — confidence 84%, second drone dispatched for confirmation',
      'Thermal imaging for hotspot and gas cloud detection across all pit faces',
      'AI path planning optimises coverage across all 9 operational zones simultaneously',
    ],
    'stat':{ 'value':'6', 'label':'Drones active in flight' }, 'route':'/drone-monitoring' },
  { 'id':13, 'type':'dashboard', 'title':'Ground Sensor Network',
    'subtitle':'89 sensors — gas, dust, temperature, vibration across 9 zones',
    'content':[
      '89 sensors monitoring CH4, CO, H2S, SO2, PM2.5, PM10, temperature, humidity, and vibration',
      'GS-11: CH4 8.4 ppm at Fuel Zone 8 — AI predicts threshold breach within 2 hours',
      'GS-08: PM2.5 58.4 ug/m3 — approaching environmental regulatory limit',
      'Zone-level heatmap with live readings; configurable alert thresholds per sensor type',
    ],
    'stat':{ 'value':'89', 'label':'Ground sensors online' }, 'route':'/ground-sensors' },
  { 'id':14, 'type':'dashboard', 'title':'Worker Health & Safety',
    'subtitle':'247 wearables — every worker monitored continuously',
    'content':[
      '247 workers monitored across 9 zones: heart rate, SpO2, body temp, fatigue score, fall detection',
      'W004: heat stress risk — fatigue 78, HR 128 bpm; intervention recommended within 15 minutes',
      'W010: fall detected — no movement for 90 seconds; rescue team dispatched with live GPS routing',
      '47 consecutive incident-free days — best performance in this site\'s recorded history',
    ],
    'stat':{ 'value':'247', 'label':'Workers monitored in real time' }, 'route':'/worker-health' },
  { 'id':15, 'type':'dashboard', 'title':'Equipment Safety',
    'subtitle':'8 heavy machines — predictive maintenance active',
    'content':[
      '8 machines: ZX870 excavators, HD785 haul trucks, A60H articulated dumpers tracked 24/7',
      'EQ08: engine temp 103C — critical alert; failure predicted within 10 hours without intervention',
      'Predictive maintenance ML models forecast failures 24–72 hours ahead',
      'Proximity alerts warn operators when vehicles approach worker exclusion zones',
    ],
    'stat':{ 'value':'8', 'label':'Heavy machines monitored' }, 'route':'/equipment-safety' },
  { 'id':16, 'type':'dashboard', 'title':'Fleet Management',
    'subtitle':'Real-time vehicle tracking and route compliance',
    'content':[
      'Live GPS tracking of all vehicles: position, speed, heading, and zone assignment',
      'Operator licence and safety training compliance — expiry alerts sent 90 days in advance',
      'AI-optimised haul routes reduce fuel consumption by 12% and cut idle time by 18%',
      'Proximity alert system: vehicle-to-vehicle and vehicle-to-worker collision prevention',
    ],
    'stat':{ 'value':'12%', 'label':'Fuel savings via AI route optimisation' }, 'route':'/fleet-management' },
  { 'id':17, 'type':'dashboard', 'title':'AI Prediction Engine',
    'subtitle':'Detect risk before it becomes an incident',
    'content':[
      'ML models for gas accumulation, equipment failure, worker fatigue, slope stability, noise, route deviation',
      'Each prediction shows confidence score, time-to-event, and top contributing sensor factors',
      'Explainable AI: every recommendation includes the reasoning chain — no black box decisions',
      'What-if simulation: test intervention actions and see simulated risk reduction before acting',
    ],
    'stat':{ 'value':'89%', 'label':'Maximum prediction confidence (Gas Safety Model)' }, 'route':'/ai-prediction' },
  { 'id':18, 'type':'dashboard', 'title':'Incident & Environmental Response',
    'subtitle':'From detection to resolution in minutes',
    'content':[
      'Automated incident workflow: Alert → Verify → Assign → Escalate → Resolve → Review',
      'AI evacuation route planner: calculates safest path in real time based on active hazards',
      'Rescue team dispatch: GPS-tracked responders reach any zone within 8 minutes',
      'Full audit trail — every AI decision logged with timestamp for post-incident compliance review',
    ],
    'stat':{ 'value':'< 8 min', 'label':'Rescue team response time' }, 'route':'/incident-response' },
  { 'id':19, 'type':'dashboard', 'title':'Sustainability Intelligence',
    'subtitle':'AI safety systems are sustainability systems',
    'content':[
      'Carbon intensity 18.4 kg CO2e/t — down 4.2% YoY via AI route and process optimisation',
      'Water recycling 76% — AI identifies processing bottlenecks pushing toward 85% target',
      'Land rehabilitation 64% complete — drone monitoring tracks revegetation progress weekly',
      'SDG 3, 8, 9, 12, 13, 15 and 16 actively tracked; ESG readiness score 87%',
    ],
    'stat':{ 'value':'87%', 'label':'ESG Compliance Readiness Score' }, 'route':'/sustainability' },
  { 'id':20, 'type':'dashboard', 'title':'Noise Exposure Intelligence',
    'subtitle':'Preventing noise-induced hearing loss with AI and wearables',
    'content':[
      '64 workers monitored for noise dose in real time using wearable dosimeters and zone beacons',
      '6 workers currently above the 85 dBA 8-hour action level — rotation and PPE alerts issued',
      'AI anomaly detection: Crusher C2 bearing noise pattern changed (78% conf.) — inspection scheduled',
      'Hierarchy-of-controls recommendations generated automatically from sensor and dose data',
    ],
    'stat':{ 'value':'91%', 'label':'Hearing PPE compliance — verified in high-noise zones' }, 'route':'/noise-exposure' },
  { 'id':22, 'type':'content', 'title':'Business Value Summary',
    'subtitle':'Safer workers. Lower impact. Better governance.',
    'content':[
      '34% reduction in safety incidents year-on-year across all monitored zones',
      '12% reduction in fleet fuel consumption via AI-optimised haul routes',
      '280 kWh saved per night via smart ventilation scheduling in underground zones',
      '$2.4M estimated annual savings from predictive maintenance alone',
      'ESG readiness score of 87% — enabling investor-grade sustainability disclosure',
    ],
    'stat':{ 'value':'$2.4M', 'label':'Estimated annual savings — predictive maintenance' },
  },
  { 'id':22, 'type':'takeaway', 'title':'Takeaway', 'icon':'✨',
    'statement':'MineSafe demonstrates how AI can move industrial safety from reactive response to predictive prevention. The system connects data, explains risk, recommends action, and supports safer, more sustainable mining operations.',
    'closing':'"The goal is not to replace human judgment, but to give safety teams earlier, clearer, and more actionable intelligence."',
  },
  { 'id':23, 'type':'thankyou', 'title':'Thank You',
    'subtitle':'Harbin Institute of Technology, Shenzhen',
    'institution':'HITSZ — AI for Sustainable Development · 2026', 'icon':'🎓',
    'members':['Harta Deddy Irawan','Sashi Kumar','Isabella Gbefa','Shania Clauren'],
  },
  { 'id':24, 'type':'videothankyou', 'title':'Thank You',
    'subtitle':'AI for Sustainable Development — HITSZ 2026',
    'closing':'From data, safety, and sustainability to smarter mining decisions.',
    'footer':'MineSafe AI System · HITSZ 2026',
    'downloadSrc':'https://raw.githubusercontent.com/hdyrawan/ai-mining-safety-dashboard/3ee1a6f/public/images/mining/Thank%20you%20video.MP4',
  },
]

# ── Build ─────────────────────────────────────────────────────────────────────

prs = new_prs()

builders = {
    'title':        slide_title,
    'team':         slide_team,
    'content':      slide_content,
    'problem':      slide_problem,
    'comparison':   slide_comparison,
    'solution':     slide_solution,
    'agents':       slide_agents,
    'training':     slide_training,
    'visits':       slide_visits,
    'dashboard':    slide_dashboard,
    'takeaway':     slide_takeaway,
    'thankyou':     slide_thankyou,
    'videothankyou':slide_video_thankyou,
}

for i, s in enumerate(slides_data):
    fn = builders.get(s['type'])
    if fn:
        fn(prs, s, i + 1)
        print(f"  [{i+1:02d}] {s['type']:16s}  {s['title']}")

out = "/home/irawanhd/ai-sustainable-dashboard/MineSafe AI System — HITSZ 2026.pptx"
prs.save(out)
print(f"\nSaved → {out}")
print(f"Slides: {len(prs.slides)}")
